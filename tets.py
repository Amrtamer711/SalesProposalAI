from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn

# === Force Visible Borders ===
def set_cell_border(cell):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    # Remove any existing border elements
    for edge in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        el = tcPr.find(qn(edge))
        if el is not None:
            tcPr.remove(el)

    # Add explicit borders
    for edge in ("L", "R", "T", "B"):
        ln = OxmlElement(f"a:ln{edge}")
        ln.set("w", "12700")  # 1pt = 12700 EMUs

        solidFill = OxmlElement("a:solidFill")
        srgbClr = OxmlElement("a:srgbClr")
        srgbClr.set("val", "000000")  # Black
        solidFill.append(srgbClr)
        ln.append(solidFill)

        prstDash = OxmlElement("a:prstDash")
        prstDash.set("val", "solid")
        ln.append(prstDash)

        tcPr.append(ln)

# === Config ===
IMAGE_PATH = "image.png"
OUTPUT_FILE = "final_proposal.pptx"

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# === Layout ===
rows = 8
left = Inches(0.65)
top = Inches(0.93)
table_width = Inches(12.0)
col1_width = Inches(3.0)
col2_width = table_width - col1_width

# === Data
data = [
    ("Location:", "The Triple Crown – 3 Digital Unipoles – 6 Screens – fully synched – 1 Spot– 16 Seconds – 16.6% SOV – Total Loop is 6 spots"),
    ("Start Date:", "1st December 2025"),
    ("Duration:", ["2 Weeks", "4 Weeks", "6 Weeks"]),
    ("Net Rate:", ["AED 1,250,000", "AED 2,300,000", "AED 3,300,000"]),
    ("Upload Fee:", "AED 3,000"),
    ("Municipality Fee:", "AED 520 Per Image/Message"),
    ("VAT 5% :", ["AED 62,676", "AED 115,000", "AED 170,000"]),
    ("Total:", ["AED 1,316,196", "AED 2,418,000", "AED 3,500,000"]),
]

split_start_index = 2
max_splits = max(len(v) if isinstance(v, list) else 1 for _, v in data[split_start_index:])
cols = 1 + max_splits

# === Header image
title_image = slide.shapes.add_picture(IMAGE_PATH, left, top, width=table_width)
row_height = title_image.height
table_height = int(row_height * rows)

# === Table
table_shape = slide.shapes.add_table(rows, cols, left, top, table_width, table_height)
table = table_shape.table

# ✅ Disable PowerPoint's default styles to show borders
table.style = None
table.first_row = False
table.horz_banding = False
table.vert_banding = False

# Column widths
table.columns[0].width = col1_width
split_col_width = int(col2_width / (cols - 1))
for j in range(1, cols):
    table.columns[j].width = split_col_width

# Row heights
for row in table.rows:
    row.height = int(table_height / rows)

# === Fill table
for i, (label, value) in enumerate(data):
    label_cell = table.cell(i, 0)

    if i == 0:
        label_cell.merge(table.cell(i, cols - 1))
        label_cell.fill.background()
        tf = label_cell.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = "Financial Proposal"
        run.font.size = Pt(24)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)
        continue

    label_cell.text = label
    label_cell.fill.solid()
    label_cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
    tf = label_cell.text_frame
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    run = tf.paragraphs[0].runs[0]
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(0, 0, 0)

    if isinstance(value, list):
        for j, val in enumerate(value):
            val_cell = table.cell(i, j + 1)
            val_cell.text = val
            val_cell.fill.solid()
            val_cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            tf = val_cell.text_frame
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            run = tf.paragraphs[0].runs[0]
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(255, 0, 0) if "Net Rate" in label else RGBColor(0, 0, 0)
    else:
        val_cell = table.cell(i, 1)
        val_cell.merge(table.cell(i, cols - 1))
        val_cell.text = value
        val_cell.fill.solid()
        val_cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
        tf = val_cell.text_frame
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        run = tf.paragraphs[0].runs[0]
        run.font.size = Pt(14)
        run.font.color.rgb = RGBColor(35, 78, 173) if "Fee" in label else RGBColor(0, 0, 0)

# === Add borders
for row in table.rows:
    for cell in row.cells:
        set_cell_border(cell)

prs.save(OUTPUT_FILE)
print(f"✅ Saved to {OUTPUT_FILE}")
