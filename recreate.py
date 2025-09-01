from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn

# === Force Visible Borders ===
def set_cell_border(cell, edges=("L", "R", "T", "B")):
    """Set black borders on specified edges of a cell"""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    # Remove existing borders
    for side in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        existing = tcPr.find(qn(side))
        if existing is not None:
            tcPr.remove(existing)

    # Define border style
    for edge in edges:
        ln = OxmlElement(f"a:ln{edge}")
        ln.set("w", "25400")  # 2 pt in EMUs (thicker for visibility)
        ln.set("cap", "flat")  # flat cap
        ln.set("cmpd", "sng")  # single line
        ln.set("algn", "ctr")  # center alignment

        # Solid fill with black color
        solidFill = OxmlElement("a:solidFill")
        srgbClr = OxmlElement("a:srgbClr")
        srgbClr.set("val", "000000")  # black
        solidFill.append(srgbClr)
        ln.append(solidFill)

        # Ensure no transparency
        prstDash = OxmlElement("a:prstDash")
        prstDash.set("val", "solid")
        ln.append(prstDash)
        
        # Head and tail ends
        headEnd = OxmlElement("a:headEnd")
        headEnd.set("type", "none")
        ln.append(headEnd)
        
        tailEnd = OxmlElement("a:tailEnd") 
        tailEnd.set("type", "none")
        ln.append(tailEnd)
        
        # Round line join
        round_join = OxmlElement("a:round")
        ln.append(round_join)

        tcPr.append(ln)

# === Config ===
IMAGE_PATH = "image.png"
OUTPUT_FILE = "final_proposal.pptx"

prs = Presentation()
# Make slide even wider and taller
prs.slide_width = Inches(20)  # Much wider slide
prs.slide_height = Inches(12)  # Much taller slide for more bottom space
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# === Layout ===
rows = 9
left = Inches(0.75)      # Slightly more margin for balance
top = Inches(0.5)        # Keep table high up
table_width = Inches(18.5)  # Much wider table to use the space
col1_width = Inches(4.0)    # Proportionally wider first column
col2_width = table_width - col1_width

# === Data
data = [
    ("Financial Proposal", None),  # Header row
    ("Location:", "The Triple Crown – 3 Digital Unipoles – 6 Screens – fully synched – 1 Spot– 16 Seconds – 16.6% SOV – Total Loop is 6 spots"),
    ("Start Date:", "1st December 2025"),
    ("Duration:", ["2 Weeks", "4 Weeks", "6 Weeks"]),
    ("Net Rate:", ["AED 1,250,000", "AED 2,300,000", "AED 3,300,000"]),
    ("Upload Fee:", "AED 3,000"),
    ("Municipality Fee:", "AED 520 Per Image/Message"),
    ("VAT 5% :", ["AED 62,676", "AED 115,000", "AED 170,000"]),
    ("Total:", ["AED 1,316,196", "AED 2,418,000", "AED 3,500,000"]),
]

split_start_index = 3  # Start splitting from "Duration" row (index 3)
max_splits = max(len(v) if isinstance(v, list) else 1 for _, v in data[split_start_index:])
cols = 1 + max_splits

# === Header image
title_image = slide.shapes.add_picture(IMAGE_PATH, left, top, width=table_width)
# Make table taller by increasing row height
row_height = Inches(0.9)  # Increased row height for taller table
table_height = int(row_height * rows)  # Total height will be about 8.1 inches

# === Table
table_shape = slide.shapes.add_table(rows, cols, left, top, table_width, table_height)
table = table_shape.table

# Set column widths
table.columns[0].width = col1_width
split_col_width = int(col2_width / (cols - 1))
for j in range(1, cols):
    table.columns[j].width = split_col_width

# Set row heights
for row in table.rows:
    row.height = int(table_height / rows)

# Fill table
for i, (label, value) in enumerate(data):
    label_cell = table.cell(i, 0)

    if i == 0:  # Financial Proposal header
        label_cell.merge(table.cell(i, cols - 1))
        label_cell.fill.background()
        tf = label_cell.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label  # Use the label which is "Financial Proposal"
        run.font.size = Pt(36)  # Bigger header text for larger table
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)
        continue

    label_cell.text = label
    label_cell.fill.solid()
    # Grey background for Total row, white for others
    if label == "Total:":
        label_cell.fill.fore_color.rgb = RGBColor(128, 128, 128)  # Grey
    else:
        label_cell.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White
    tf = label_cell.text_frame
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER  # Center align text
    run = p.runs[0]
    run.font.size = Pt(20)  # Bigger text for label cells
    # White text for Total row, black for others
    if label == "Total:":
        run.font.color.rgb = RGBColor(255, 255, 255)  # White text
        run.font.bold = True  # Make it bold
        run.font.size = Pt(24)  # Even bigger for Total row
    else:
        run.font.color.rgb = RGBColor(0, 0, 0)  # Black text

    if isinstance(value, list):
        for j, val in enumerate(value):
            val_cell = table.cell(i, j + 1)
            val_cell.text = val
            val_cell.fill.solid()
            # Grey background for Total row, white for others
            if label == "Total:":
                val_cell.fill.fore_color.rgb = RGBColor(128, 128, 128)  # Grey
            else:
                val_cell.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White
            tf = val_cell.text_frame
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER  # Center align text
            run = p.runs[0]
            run.font.size = Pt(20)  # Bigger text for value cells
            # Special formatting for different rows
            if label == "Total:":
                run.font.color.rgb = RGBColor(255, 255, 255)  # White text
                run.font.bold = True  # Make it bold
                run.font.size = Pt(24)  # Even bigger for Total row
            elif "Net Rate" in label:
                run.font.color.rgb = RGBColor(255, 0, 0)  # Red text
            else:
                run.font.color.rgb = RGBColor(0, 0, 0)  # Black text
    else:
        val_cell = table.cell(i, 1)
        val_cell.merge(table.cell(i, cols - 1))
        val_cell.text = value
        val_cell.fill.solid()
        # Grey background for Total row, white for others
        if label == "Total:":
            val_cell.fill.fore_color.rgb = RGBColor(128, 128, 128)  # Grey
        else:
            val_cell.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White
        tf = val_cell.text_frame
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER  # Center align text
        run = p.runs[0]
        run.font.size = Pt(20)  # Bigger text for merged cells
        # Special formatting for different rows
        if label == "Total:":
            run.font.color.rgb = RGBColor(255, 255, 255)  # White text
            run.font.bold = True  # Make it bold
            run.font.size = Pt(24)  # Even bigger for Total row
        elif "Fee" in label:
            run.font.color.rgb = RGBColor(35, 78, 173)  # Blue text
        else:
            run.font.color.rgb = RGBColor(0, 0, 0)  # Black text

# === Add borders visibly
# First, set borders on all cells
for row in table.rows:
    for cell in row.cells:
        set_cell_border(cell)

# Additionally, ensure table has default borders
table_element = table._tbl
tblPr = table_element.find(qn('a:tblPr'))
if tblPr is None:
    tblPr = OxmlElement('a:tblPr')
    table_element.insert(0, tblPr)

# Remove any existing table style that might override borders
for style in tblPr.findall(qn('a:tableStyleId')):
    tblPr.remove(style)

# === Add bullet points below table ===
bullet_text = """• A DM fee of AED 520 per image/message applies. The final fee will be confirmed after the final artwork is received.
• An official booking order is required to secure the location/spot.
• Once a booking is confirmed, cancellations are not allowed even in case an artwork is rejected by the authorities, the client will be required to submit a revised artwork.
• All artworks are subject to approval by BackLite Media and DM.
• Location availability is subject to change.
• The artwork must comply with DM's guidelines.
• This proposal is valid until the 24th of November, 2025."""

# Add text box for bullet points
bullet_box = slide.shapes.add_textbox(
    left=Inches(0.75),
    top=Inches(9.0),  # Position below the taller table
    width=Inches(18.5),
    height=Inches(2.5)
)

tf = bullet_box.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0)
tf.margin_right = Inches(0)
tf.margin_top = Inches(0.1)
tf.margin_bottom = Inches(0)

p = tf.paragraphs[0]
p.text = bullet_text
p.font.size = Pt(14)  # Bigger text (was 12)
p.font.color.rgb = RGBColor(0, 0, 0)
p.line_spacing = 1.3  # Slightly more line spacing for readability

# === Space calculation ===
# Table starts at 0.5" and is about 8.1" tall, ending at ~8.6"
# Slide height is 12", leaving about 3.4" of space at the bottom for text

prs.save(OUTPUT_FILE)
print(f"✅ Saved to {OUTPUT_FILE}")
print(f"📐 Table dimensions: {table_width/Inches(1):.1f}\" wide x {table_height/Inches(1):.1f}\" tall")
print(f"📍 Space below table: ~{(prs.slide_height - top - Inches(table_height/Inches(1)))/Inches(1):.1f}\" for additional text")
