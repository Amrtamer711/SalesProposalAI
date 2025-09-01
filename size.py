from pptx import Presentation
from pptx.util import Inches
import os

# === Input ===
SOURCE_PPTX = "Template.pptx"  # replace with your filename

# === Load presentation ===
prs = Presentation(SOURCE_PPTX)

# === Slide size ===
slide_width_in = prs.slide_width / 914400
slide_height_in = prs.slide_height / 914400
print(f"📐 Slide size: {slide_width_in:.2f}\" x {slide_height_in:.2f}\"")

# === Look for a table shape in the first slide ===
first_slide = prs.slides[0]
table_found = False

for shape in first_slide.shapes:
    if shape.has_table:
        table_found = True
        left_in = shape.left / 914400
        top_in = shape.top / 914400
        width_in = shape.width / 914400
        height_in = shape.height / 914400

        print("🗂️  Table found:")
        print(f"   • Position: Left = {left_in:.2f}\", Top = {top_in:.2f}\"")
        print(f"   • Size: Width = {width_in:.2f}\", Height = {height_in:.2f}\"")
        break

if not table_found:
    print("❌ No table found in the first slide.")
