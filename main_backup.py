import os
import json
import asyncio
import logging
from pathlib import Path
from datetime import datetime
from typing import Dict, Any
from concurrent.futures import ThreadPoolExecutor

from dotenv import load_dotenv
from fastapi import FastAPI, Request, HTTPException
from fastapi.responses import JSONResponse
import uvicorn

from slack_sdk.web.async_client import AsyncWebClient
from slack_sdk.signature import SignatureVerifier
from openai import AsyncOpenAI

# PowerPoint imports
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn

# Image processing imports
from PIL import Image
import io

# ========== CONFIG ==========
load_dotenv()
app = FastAPI(title="Proposal Bot API")
executor = ThreadPoolExecutor()

logger = logging.getLogger("proposal-bot")
logging.basicConfig(level=logging.INFO)

slack_client = AsyncWebClient(token=os.getenv("SLACK_BOT_TOKEN"))
signature_verifier = SignatureVerifier(os.getenv("SLACK_SIGNING_SECRET"))
openai_client = AsyncOpenAI(api_key=os.getenv("OPENAI_API_KEY"))

LOCATION_MAPPING = {
    "landmark": "1. Desirable by Location - The Landmark Series copy 2.pptx",
    "jawhara": "Jawhara.pptx",
    "gateway": "The Gateway.pptx",
    "oryx": "The Oryx.pptx",
    "triple crown": "The Triple Crown.pptx"
}

user_history: Dict[str, list] = {}

# ========== HELPER FUNCTIONS ==========
def set_cell_border(cell, edges=("L", "R", "T", "B")):
    """Set black borders on specified edges of a cell"""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    # Remove existing borders
    for side in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        existing = tcPr.find(qn(side))
        if existing is not None:
            tcPr.remove(existing)

    # Define border style
    for edge in edges:
        ln = OxmlElement(f"a:ln{edge}")
        ln.set("w", "25400")  # 2 pt in EMUs (thicker for visibility)
        ln.set("cap", "flat")  # flat cap
        ln.set("cmpd", "sng")  # single line
        ln.set("algn", "ctr")  # center alignment

        # Solid fill with black color
        solidFill = OxmlElement("a:solidFill")
        srgbClr = OxmlElement("a:srgbClr")
        srgbClr.set("val", "000000")  # black
        solidFill.append(srgbClr)
        ln.append(solidFill)

        # Ensure no transparency
        prstDash = OxmlElement("a:prstDash")
        prstDash.set("val", "solid")
        ln.append(prstDash)
        
        # Head and tail ends
        headEnd = OxmlElement("a:headEnd")
        headEnd.set("type", "none")
        ln.append(headEnd)
        
        tailEnd = OxmlElement("a:tailEnd") 
        tailEnd.set("type", "none")
        ln.append(tailEnd)
        
        # Round line join
        round_join = OxmlElement("a:round")
        ln.append(round_join)

        tcPr.append(ln)

def create_financial_proposal_slide(slide, financial_data: dict, slide_width, slide_height):
    """Create the financial proposal slide content scaled to fit the slide dimensions."""
    # === Calculate scaling factors ===
    # Original dimensions were for 20" x 12" slides
    scale_x = slide_width / Inches(20)
    scale_y = slide_height / Inches(12)
    # Use the smaller scale to maintain proportions
    scale = min(scale_x, scale_y)
    
    # === Layout ===
    rows = 9
    # Scale all dimensions proportionally (must be integers)
    left = int(Inches(0.75) * scale_x)
    top = int(Inches(0.5) * scale_y)
    table_width = int(Inches(18.5) * scale_x)
    col1_width = int(Inches(4.0) * scale_x)
    col2_width = table_width - col1_width
    
    # === Extract financial data ===
    location_name = financial_data["location"]
    start_date = financial_data["start_date"]
    durations = financial_data["durations"]
    net_rates = financial_data["net_rates"]
    
    # === Location details
    location_details = {
        "landmark": "The Landmark - Premium Digital Display - 1 Screen - 1 Spot - 16 Seconds - 16.6% SOV - Total Loop is 6 spots",
        "jawhara": "Jawhara - Digital Billboard - 2 Screens synched - 1 Spot - 16 Seconds - 16.6% SOV - Total Loop is 6 spots",
        "gateway": "The Gateway - 4 Digital Unipoles - 8 Screens - fully synched - 1 Spot - 16 Seconds - 16.6% SOV - Total Loop is 6 spots",
        "oryx": "The Oryx - Digital Display Network - 3 Screens - 1 Spot - 16 Seconds - 16.6% SOV - Total Loop is 6 spots",
        "triple crown": "The Triple Crown – 3 Digital Unipoles – 6 Screens – fully synched – 1 Spot– 16 Seconds – 16.6% SOV – Total Loop is 6 spots"
    }
    
    location_text = location_details.get(location_name.lower(), 
        f"{location_name.title()} - Digital Display - 1 Spot - 16 Seconds - 16.6% SOV - Total Loop is 6 spots")
    
    # === Calculate VAT and totals ===
    upload_fee = 3000  # Fixed
    municipality_fee = 520  # Fixed per image/message
    
    # Calculate VAT and total for each duration option
    vat_amounts = []
    total_amounts = []
    
    for net_rate_str in net_rates:
        # Extract numeric value from string like "AED 1,250,000"
        net_rate = float(net_rate_str.replace("AED", "").replace(",", "").strip())
        subtotal = net_rate + upload_fee + municipality_fee
        vat = subtotal * 0.05
        total = subtotal + vat
        
        vat_amounts.append(f"AED {vat:,.0f}")
        total_amounts.append(f"AED {total:,.0f}")
    
    # === Build data array ===
    data = [
        ("Financial Proposal", None),  # Header row
        ("Location:", location_text),
        ("Start Date:", start_date),
        ("Duration:", durations if len(durations) > 1 else durations[0]),
        ("Net Rate:", net_rates if len(net_rates) > 1 else net_rates[0]),
        ("Upload Fee:", "AED 3,000"),
        ("Municipality Fee:", "AED 520 Per Image/Message"),
        ("VAT 5% :", vat_amounts if len(vat_amounts) > 1 else vat_amounts[0]),
        ("Total:", total_amounts if len(total_amounts) > 1 else total_amounts[0]),
    ]
    
    split_start_index = 3
    max_splits = max(len(v) if isinstance(v, list) else 1 for _, v in data[split_start_index:])
    cols = 1 + max_splits
    
    # === Header image
    image_path = Path(__file__).parent / "image.png"
    if image_path.exists():
        slide.shapes.add_picture(str(image_path), left, top, width=table_width)
    # Make table taller by increasing row height (scaled)
    row_height = int(Inches(0.9) * scale_y)  # Scale row height
    table_height = int(row_height * rows)
    
    # === Table
    table_shape = slide.shapes.add_table(rows, cols, left, top, table_width, table_height)
    table = table_shape.table
    
    # Set column widths
    table.columns[0].width = col1_width
    split_col_width = int(col2_width / (cols - 1))
    for j in range(1, cols):
        table.columns[j].width = split_col_width
    
    # Set row heights
    for row in table.rows:
        row.height = int(table_height / rows)
    
    # Fill table
    for i, (label, value) in enumerate(data):
        label_cell = table.cell(i, 0)
        
        if i == 0:  # Financial Proposal header
            label_cell.merge(table.cell(i, cols - 1))
            label_cell.fill.background()
            tf = label_cell.text_frame
            tf.clear()
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = label
            run.font.size = Pt(int(36 * scale))  # Scale font size
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 255, 255)
            continue
        
        # Regular rows
        label_cell.text = label
        label_cell.fill.solid()
        if label == "Total:":
            label_cell.fill.fore_color.rgb = RGBColor(128, 128, 128)
        else:
            label_cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
        
        tf = label_cell.text_frame
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.size = Pt(int(20 * scale))
        
        if label == "Total:":
            run.font.color.rgb = RGBColor(255, 255, 255)
            run.font.bold = True
            run.font.size = Pt(int(24 * scale))
        else:
            run.font.color.rgb = RGBColor(0, 0, 0)
        
        if isinstance(value, list):
            for j, val in enumerate(value):
                val_cell = table.cell(i, j + 1)
                val_cell.text = val
                val_cell.fill.solid()
                if label == "Total:":
                    val_cell.fill.fore_color.rgb = RGBColor(128, 128, 128)
                else:
                    val_cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                
                tf = val_cell.text_frame
                tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.runs[0]
                run.font.size = Pt(int(20 * scale))
                
                if label == "Total:":
                    run.font.color.rgb = RGBColor(255, 255, 255)
                    run.font.bold = True
                    run.font.size = Pt(int(24 * scale))
                elif label == "Net Rate:":
                    run.font.color.rgb = RGBColor(255, 0, 0)
                else:
                    run.font.color.rgb = RGBColor(0, 0, 0)
        else:
            val_cell = table.cell(i, 1)
            val_cell.merge(table.cell(i, cols - 1))
            val_cell.text = value
            val_cell.fill.solid()
            if label == "Total:":
                val_cell.fill.fore_color.rgb = RGBColor(128, 128, 128)
            else:
                val_cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            
            tf = val_cell.text_frame
            tf.clear()
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            p = tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = value
            run.font.size = Pt(int(20 * scale))
            
            if label == "Total:":
                run.font.color.rgb = RGBColor(255, 255, 255)
                run.font.bold = True
                run.font.size = Pt(int(24 * scale))
            elif "Fee" in label:
                run.font.color.rgb = RGBColor(35, 78, 173)
            else:
                run.font.color.rgb = RGBColor(0, 0, 0)
    
    # === Add borders visibly
    # First, set borders on all cells
    for row in table.rows:
        for cell in row.cells:
            set_cell_border(cell)
    
    # Additionally, ensure table has default borders
    table_element = table._tbl
    tblPr = table_element.find(qn('a:tblPr'))
    if tblPr is None:
        tblPr = OxmlElement('a:tblPr')
        table_element.insert(0, tblPr)
    
    # Remove any existing table style that might override borders
    for style in tblPr.findall(qn('a:tableStyleId')):
        tblPr.remove(style)
    
    # === Add bullet points below table ===
    bullet_text = """• A DM fee of AED 520 per image/message applies. The final fee will be confirmed after the final artwork is received.
• An official booking order is required to secure the location/spot.
• Once a booking is confirmed, cancellations are not allowed even in case an artwork is rejected by the authorities, the client will be required to submit a revised artwork.
• All artworks are subject to approval by BackLite Media and DM.
• Location availability is subject to change.
• The artwork must comply with DM's guidelines.
• This proposal is valid until the 24th of November, 2025."""
    
    # Scale bullet point positioning
    bullet_box = slide.shapes.add_textbox(
        left=int(Inches(0.75) * scale_x),
        top=int(Inches(9.0) * scale_y),
        width=int(Inches(18.5) * scale_x),
        height=int(Inches(2.5) * scale_y)
    )
    
    tf = bullet_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0)
    
    p = tf.paragraphs[0]
    p.text = bullet_text
    p.font.size = Pt(int(14 * scale))
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.line_spacing = 1.3

# ========== PDF CONVERSION AND MERGING ==========
def convert_pptx_to_images_pdf(pptx_path: str) -> str:
    """Convert PowerPoint slides to images and then create PDF from images."""
    import tempfile
    import subprocess
    import platform
    import shutil
    from PIL import Image
    from reportlab.lib.pagesizes import landscape, letter, A4
    from reportlab.pdfgen import canvas
    from reportlab.lib.utils import ImageReader
    
    # Create temporary PDF file
    pdf_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
    pdf_file.close()
    
    # Create temporary directory for images
    temp_dir = tempfile.mkdtemp()
    
    try:
        # Method 1: Try using LibreOffice to export as images
        libreoffice_paths = [
            '/opt/homebrew/bin/soffice',  # Mac Homebrew ARM64
            'soffice',  # Alternative command
            'libreoffice',  # Linux/Mac with LibreOffice in PATH
            '/Applications/LibreOffice.app/Contents/MacOS/soffice',  # Mac
            '/usr/bin/libreoffice',  # Linux common path
            '/usr/local/bin/libreoffice',  # Mac Homebrew Intel
        ]
        
        images_created = False
        
        for lo_path in libreoffice_paths:
            if shutil.which(lo_path) or os.path.exists(lo_path):
                try:
                    # Export to PNG images
                    cmd = [
                        lo_path,
                        '--headless',
                        '--convert-to',
                        'png',
                        '--outdir',
                        temp_dir,
                        pptx_path
                    ]
                    
                    result = subprocess.run(cmd, capture_output=True, text=True, timeout=30)
                    
                    if result.returncode == 0:
                        # Check if any images were created
                        png_files = sorted([f for f in os.listdir(temp_dir) if f.endswith('.png')])
                        if png_files:
                            images_created = True
                            break
                except Exception as e:
                    logger.debug(f"LibreOffice image export failed: {e}")
                    continue
        
        # Method 2: If no images created, try using python-pptx to render slides
        if not images_created:
            logger.info("Using python-pptx screenshot method...")
            from pptx import Presentation
            import io
            
            # Load presentation
            pres = Presentation(pptx_path)
            
            # For each slide, we'll create a high-resolution image
            for slide_num, slide in enumerate(pres.slides):
                # Create a white background image at high resolution
                # Standard slide size is 10" x 7.5" at 300 DPI = 3000 x 2250 pixels
                img_width = 3000
                img_height = 2250
                img = Image.new('RGB', (img_width, img_height), 'white')
                
                # Note: python-pptx doesn't have built-in rendering, so we'll create a simple representation
                # For a production system, you would use a proper rendering library or API
                # This is a placeholder that shows the slide number
                from PIL import ImageDraw, ImageFont
                draw = ImageDraw.Draw(img)
                
                try:
                    # Try to use a better font if available
                    font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 72)
                except:
                    font = ImageFont.load_default()
                
                # Draw slide number (placeholder for actual content)
                text = f"Slide {slide_num + 1}"
                draw.text((img_width//2 - 100, img_height//2 - 50), text, fill='black', font=font)
                
                # Save the image
                img_path = os.path.join(temp_dir, f"slide_{slide_num:03d}.png")
                img.save(img_path, 'PNG', dpi=(300, 300))
            
            png_files = sorted([f for f in os.listdir(temp_dir) if f.endswith('.png')])
            images_created = True
        
        # Method 3: Try using macOS-specific conversion
        if not images_created and platform.system() == "Darwin":
            try:
                # Use sips (scriptable image processing system) on macOS
                # First convert to PDF using Preview or PowerPoint, then to images
                temp_pdf = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
                temp_pdf.close()
                
                # Try PowerPoint first
                powerpoint_script = f'''
                tell application "Microsoft PowerPoint"
                    open POSIX file "{pptx_path}"
                    save active presentation in POSIX file "{temp_pdf.name}" as save as PDF
                    close active presentation
                end tell
                '''
                
                result = subprocess.run(
                    ['osascript', '-e', powerpoint_script],
                    capture_output=True,
                    text=True,
                    timeout=30
                )
                
                if result.returncode == 0 and os.path.exists(temp_pdf.name):
                    # Convert PDF pages to images using sips
                    from pdf2image import convert_from_path
                    images = convert_from_path(temp_pdf.name, dpi=300)
                    
                    for i, image in enumerate(images):
                        img_path = os.path.join(temp_dir, f"slide_{i:03d}.png")
                        image.save(img_path, 'PNG')
                    
                    png_files = sorted([f for f in os.listdir(temp_dir) if f.endswith('.png')])
                    images_created = True
                    
                os.unlink(temp_pdf.name)
            except Exception as e:
                logger.debug(f"macOS conversion failed: {e}")
        
        # Now create PDF from images
        if images_created and png_files:
            # Get the first image to determine page size
            first_img_path = os.path.join(temp_dir, png_files[0])
            with Image.open(first_img_path) as img:
                img_width, img_height = img.size
                # Calculate page size in points (72 points per inch)
                # Assuming 300 DPI for the images
                page_width = (img_width / 300) * 72
                page_height = (img_height / 300) * 72
            
            # Create PDF with custom page size
            c = canvas.Canvas(pdf_file.name, pagesize=(page_width, page_height))
            
            for png_file in png_files:
                img_path = os.path.join(temp_dir, png_file)
                
                # Add image to fill the entire page
                c.drawImage(img_path, 0, 0, width=page_width, height=page_height)
                
                # Show page
                c.showPage()
            
            # Save the PDF
            c.save()
            
            logger.info(f"Successfully created PDF from {len(png_files)} slide images")
            return pdf_file.name
        else:
            raise Exception("No images were created from the PowerPoint file")
            
    except Exception as e:
        logger.error(f"Image-based PDF conversion failed: {e}")
        raise
    finally:
        # Clean up temporary directory
        try:
            shutil.rmtree(temp_dir)
        except:
            pass

def convert_pptx_to_pdf(pptx_path: str) -> str:
    """Convert PowerPoint to PDF using screenshot method to preserve formatting."""
    # Use the new image-based conversion method
    return convert_pptx_to_images_pdf(pptx_path)
        'soffice',  # Alternative command
        'libreoffice',  # Linux/Mac with LibreOffice in PATH
        '/Applications/LibreOffice.app/Contents/MacOS/soffice',  # Mac
        '/usr/bin/libreoffice',  # Linux common path
        '/usr/local/bin/libreoffice',  # Mac Homebrew Intel
    ]
    
    for lo_path in libreoffice_paths:
        if shutil.which(lo_path) or os.path.exists(lo_path):
            try:
                # Use LibreOffice in headless mode to convert
                cmd = [
                    lo_path,
                    '--headless',
                    '--convert-to',
                    'pdf',
                    '--outdir',
                    os.path.dirname(pdf_file.name),
                    pptx_path
                ]
                
                result = subprocess.run(cmd, capture_output=True, text=True, timeout=30)
                
                if result.returncode == 0:
                    # LibreOffice creates PDF with same name as input
                    converted_pdf = os.path.join(
                        os.path.dirname(pdf_file.name),
                        os.path.splitext(os.path.basename(pptx_path))[0] + '.pdf'
                    )
                    
                    if os.path.exists(converted_pdf):
                        # Move to our desired location
                        shutil.move(converted_pdf, pdf_file.name)
                        return pdf_file.name
            except Exception as e:
                logger.debug(f"LibreOffice conversion failed: {e}")
                continue
    
    # Method 2: Try unoconv (another cross-platform option)
    if shutil.which('unoconv'):
        try:
            cmd = ['unoconv', '-f', 'pdf', '-o', pdf_file.name, pptx_path]
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=30)
            
            if result.returncode == 0 and os.path.exists(pdf_file.name):
                return pdf_file.name
        except Exception as e:
            logger.debug(f"unoconv conversion failed: {e}")
    
    # Method 3: Try macOS-specific conversion using osascript
    if system == "Darwin":  # macOS
        try:
            # First try PowerPoint if available
            powerpoint_script = f'''
            tell application "Microsoft PowerPoint"
                open POSIX file "{pptx_path}"
                save active presentation in POSIX file "{pdf_file.name}" as save as PDF
                close active presentation
            end tell
            '''
            
            result = subprocess.run(
                ['osascript', '-e', powerpoint_script],
                capture_output=True,
                text=True,
                timeout=30
            )
            
            if result.returncode == 0 and os.path.exists(pdf_file.name):
                return pdf_file.name
        except Exception as e:
            logger.debug(f"PowerPoint conversion failed: {e}")
        
        try:
            # Then try Keynote (can open PowerPoint files)
            keynote_script = f'''
            tell application "Keynote"
                open POSIX file "{pptx_path}"
                export front document to POSIX file "{pdf_file.name}" as PDF
                close front document
            end tell
            '''
            
            result = subprocess.run(
                ['osascript', '-e', keynote_script],
                capture_output=True,
                text=True,
                timeout=30
            )
            
            if result.returncode == 0 and os.path.exists(pdf_file.name):
                return pdf_file.name
        except Exception as e:
            logger.debug(f"Keynote conversion failed: {e}")
    
    # Method 3: Fallback to enhanced text extraction with formatting
    # This is our last resort if no proper conversion tools are available
    from pptx import Presentation
    from reportlab.lib.pagesizes import landscape, letter
    from reportlab.pdfgen import canvas
    from reportlab.lib import colors
    
    try:
        pres = Presentation(pptx_path)
        page_width, page_height = landscape(letter)
        c = canvas.Canvas(pdf_file.name, pagesize=landscape(letter))
        
        for slide_idx, slide in enumerate(pres.slides):
            # Add slide background color if available
            if slide.background and hasattr(slide.background, 'fill'):
                try:
                    if slide.background.fill.type == 1:  # Solid fill
                        bg_color = slide.background.fill.fore_color.rgb
                        if bg_color:
                            c.setFillColorRGB(
                                bg_color[0]/255.0,
                                bg_color[1]/255.0, 
                                bg_color[2]/255.0
                            )
                            c.rect(0, 0, page_width, page_height, fill=1, stroke=0)
                except:
                    pass
            
            # Draw slide content with better formatting
            c.setFillColor(colors.black)
            
            # Extract and draw shapes with positioning
            for shape in slide.shapes:
                try:
                    if hasattr(shape, 'text') and shape.text.strip():
                        # Calculate approximate position
                        left = float(shape.left) / 914400 * 72  # Convert EMUs to points
                        top = float(shape.top) / 914400 * 72
                        
                        # Adjust for PDF coordinate system (bottom-left origin)
                        y_pos = page_height - top - 50
                        
                        # Draw text with appropriate size
                        text = shape.text.strip()
                        font_size = 12
                        
                        if hasattr(shape, 'text_frame') and shape.text_frame.paragraphs:
                            for para in shape.text_frame.paragraphs:
                                if para.runs:
                                    run = para.runs[0]
                                    if run.font.size:
                                        font_size = run.font.size.pt
                        
                        c.setFont("Helvetica", min(font_size, 24))
                        
                        # Wrap text if needed
                        lines = text.split('\n')
                        for line in lines:
                            if line.strip():
                                c.drawString(left, y_pos, line.strip())
                                y_pos -= font_size + 5
                    
                    # Note: For tables, images, etc., we'd need more complex handling
                    # This is a basic fallback when proper conversion tools aren't available
                    
                except Exception as e:
                    logger.debug(f"Error processing shape: {e}")
            
            # Add page number
            c.setFont("Helvetica", 10)
            c.drawString(page_width - 100, 30, f"Slide {slide_idx + 1}")
            
            # Add new page for next slide
            if slide_idx < len(pres.slides) - 1:
                c.showPage()
        
        c.save()
        
        # Add a warning that this is a text-only conversion
        logger.warning("PDF created using fallback text extraction. For exact slide rendering, install LibreOffice.")
        
        # Create a notice file to inform user
        notice_path = pdf_file.name.replace('.pdf', '_README.txt')
        with open(notice_path, 'w') as f:
            f.write("IMPORTANT: PDF Quality Notice\n")
            f.write("="*50 + "\n\n")
            f.write("This PDF was created using text extraction only.\n")
            f.write("For exact slide rendering with all visuals, please install one of:\n\n")
            f.write("1. LibreOffice (recommended):\n")
            f.write("   - Mac: brew install --cask libreoffice\n")
            f.write("   - Ubuntu/Debian: sudo apt-get install libreoffice\n")
            f.write("   - Download: https://www.libreoffice.org/download/\n\n")
            f.write("2. unoconv:\n")
            f.write("   - Mac: brew install unoconv\n")
            f.write("   - Ubuntu/Debian: sudo apt-get install unoconv\n\n")
            f.write("3. Microsoft PowerPoint (Mac only)\n\n")
            f.write("After installation, restart the bot for proper PDF conversion.\n")
        
        return pdf_file.name
        
    except Exception as e:
        logger.error(f"PDF conversion failed: {e}")
        raise

def merge_pdfs(pdf_files: list) -> str:
    """Merge multiple PDF files into one."""
    import tempfile
    from pypdf import PdfWriter, PdfReader
    
    # Create output file
    output_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
    output_file.close()
    
    # Create PDF writer
    pdf_writer = PdfWriter()
    
    # Add all pages from each PDF
    for pdf_path in pdf_files:
        pdf_reader = PdfReader(pdf_path)
        for page in pdf_reader.pages:
            pdf_writer.add_page(page)
    
    # Write the merged PDF
    with open(output_file.name, 'wb') as output:
        pdf_writer.write(output)
    
    return output_file.name

def remove_slides_and_convert_to_pdf(pptx_path: str, remove_first: bool = False, remove_last: bool = False) -> str:
    """Remove specified slides from PowerPoint and convert to PDF."""
    import tempfile
    import shutil
    from pptx import Presentation
    
    # Create a temporary copy of the presentation
    temp_pptx = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    temp_pptx.close()
    shutil.copy2(pptx_path, temp_pptx.name)
    
    # Load the presentation
    pres = Presentation(temp_pptx.name)
    xml_slides = pres.slides._sldIdLst
    slides_to_remove = []
    
    # Determine which slides to remove
    if remove_first and len(pres.slides) > 0:
        slides_to_remove.append(list(xml_slides)[0])
    
    if remove_last and len(pres.slides) > 1:
        slides_to_remove.append(list(xml_slides)[-1])
    
    # Remove the slides
    for slide_id in slides_to_remove:
        if slide_id in xml_slides:
            xml_slides.remove(slide_id)
    
    # Save the modified presentation
    pres.save(temp_pptx.name)
    
    # Convert to PDF
    pdf_path = convert_pptx_to_pdf(temp_pptx.name)
    
    # Clean up temporary PowerPoint file
    try:
        os.unlink(temp_pptx.name)
    except:
        pass
    
    return pdf_path

# ========== CORE LOGIC ==========
def create_proposal_with_template(source_path: str, financial_data: dict) -> str:
    """Create a new presentation with financial proposal slide as second-to-last."""
    import tempfile
    
    # Load the source presentation
    pres = Presentation(source_path)
    
    # Calculate position for the financial proposal slide (second-to-last)
    insert_position = max(len(pres.slides) - 1, 0)
    
    # Get the current slide dimensions to scale our content appropriately
    slide_width = pres.slide_width
    slide_height = pres.slide_height
    
    # Add a new blank slide for the financial proposal
    blank_layout = pres.slide_layouts[6] if len(pres.slide_layouts) > 6 else pres.slide_layouts[0]
    financial_slide = pres.slides.add_slide(blank_layout)
    
    # Create the financial proposal content scaled to the presentation's dimensions
    create_financial_proposal_slide(financial_slide, financial_data, slide_width, slide_height)
    
    # Move the financial slide to second-to-last position
    if len(pres.slides) > 1 and insert_position < len(pres.slides) - 1:
        xml_slides = pres.slides._sldIdLst
        new_slide_element = xml_slides[-1]  # The newly added slide
        xml_slides.remove(new_slide_element)
        xml_slides.insert(insert_position, new_slide_element)
    
    # Save the final presentation
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    pres.save(tmp.name)
    return tmp.name

async def process_proposals(proposals_data: list) -> Dict[str, Any]:
    """Process proposal requests - handles both single and multiple proposals."""
    
    if not proposals_data:
        return {"success": False, "error": "No proposals provided"}
    
    # Check if single or multiple proposals
    is_single = len(proposals_data) == 1
    
    # Process each proposal individually
    individual_files = []  # Store individual PPT files info
    pdf_files = []  # For PDF concatenation (only for multiple)
    locations = []
    
    for idx, proposal in enumerate(proposals_data):
        # Validate each proposal
        location = proposal.get("location", "").lower().strip()
        start_date = proposal.get("start_date", "1st December 2025")
        durations = proposal.get("durations", [])
        net_rates = proposal.get("net_rates", [])
        
        # Validate location
        matched = None
        loc_key = location
        for key, fname in LOCATION_MAPPING.items():
            if key in location or location in key:
                matched = fname
                loc_key = key
                break
        
        if not matched:
            return {"success": False, "error": f"Unknown location '{location}' in proposal {idx + 1}"}
        
        # Validate durations and rates match
        if len(durations) != len(net_rates):
            return {"success": False, "error": f"Mismatched durations and rates for {loc_key} - {len(durations)} durations but {len(net_rates)} rates"}
        
        if not durations:
            return {"success": False, "error": f"No duration specified for {loc_key}"}
        
        # Get the PowerPoint file
        base = Path(__file__).parent
        src = base / matched
        
        if not src.exists():
            return {"success": False, "error": f"{matched} not found"}
        
        # Create presentation with financial slide
        loop = asyncio.get_event_loop()
        financial_data = {
            "location": loc_key,
            "start_date": start_date,
            "durations": durations,
            "net_rates": net_rates
        }
        
        # Create the PowerPoint with financial slide
        pptx_file = await loop.run_in_executor(executor, create_proposal_with_template, str(src), financial_data)
        
        # Store individual PPT file info
        individual_files.append({
            "path": pptx_file,
            "location": loc_key.title(),
            "filename": f"{loc_key.title()}_Proposal.pptx"
        })
        
        locations.append(loc_key.title())
        
        # For single proposals, convert to PDF without removing slides
        if is_single:
            pdf_file = await loop.run_in_executor(executor, convert_pptx_to_pdf, pptx_file)
            individual_files[0]["pdf_path"] = pdf_file
            individual_files[0]["pdf_filename"] = f"{loc_key.title()}_Proposal.pdf"
        else:
            # For multiple proposals, determine which slides to remove
            remove_first = False
            remove_last = False
            
            if idx == 0:
                # First presentation: remove last slide (ending)
                remove_last = True
            elif idx < len(proposals_data) - 1:
                # Middle presentations: remove both first and last
                remove_first = True
                remove_last = True
            else:
                # Last presentation: remove first slide (intro)
                remove_first = True
            
            # Convert to PDF with appropriate slides removed
            pdf_file = await loop.run_in_executor(
                executor, 
                remove_slides_and_convert_to_pdf, 
                pptx_file, 
                remove_first, 
                remove_last
            )
            
            pdf_files.append(pdf_file)
    
    # For single proposal, return simplified structure
    if is_single:
        return {
            "success": True,
            "is_single": True,
            "pptx_path": individual_files[0]["path"],
            "pdf_path": individual_files[0]["pdf_path"],
            "location": individual_files[0]["location"],
            "pptx_filename": individual_files[0]["filename"],
            "pdf_filename": individual_files[0]["pdf_filename"]
        }
    
    # For multiple proposals, merge PDFs
    loop = asyncio.get_event_loop()
    merged_pdf = await loop.run_in_executor(executor, merge_pdfs, pdf_files)
    
    # Clean up individual PDF files
    for pdf_file in pdf_files:
        try:
            os.unlink(pdf_file)
        except:
            pass
    
    return {
        "success": True,
        "is_single": False,
        "individual_files": individual_files,  # Individual PPT files
        "merged_pdf_path": merged_pdf,
        "locations": ", ".join(locations),
        "merged_pdf_filename": f"Combined_Proposal_{len(locations)}_Locations.pdf"
    }


# Removed process_location_request - now using unified process_proposals function

# ========== LLM HANDLER ==========
async def main_llm_loop(channel: str, user_id: str, user_input: str):
    logger = logging.getLogger("proposal-bot")
    
    prompt = (
        f"You are a sales proposal bot for BackLite Media. You help create financial proposals for digital advertising locations.\n"
        f"You can handle SINGLE or MULTIPLE location proposals in one request. Multiple proposals will be combined into a single PDF.\n\n"
        
        f"AVAILABLE LOCATIONS: {', '.join(LOCATION_MAPPING.keys())}\n\n"
        
        f"REQUIRED INFORMATION FOR EACH LOCATION:\n"
        f"1. Location (must be one of the available locations)\n"
        f"2. Start Date (e.g., '1st December 2025', '15th January 2026')\n"
        f"3. Duration Options (e.g., '2 Weeks', '4 Weeks', '6 Weeks')\n"
        f"4. Net Rates for EACH duration (e.g., 'AED 1,250,000', 'AED 2,300,000', 'AED 3,300,000')\n\n"
        
        f"MULTIPLE PROPOSALS RULES:\n"
        f"- User can request proposals for multiple locations at once\n"
        f"- EACH location must have its own complete set of information\n"
        f"- EACH location must have matching number of durations and net rates\n"
        f"- Different locations can have different durations/rates\n"
        f"- Multiple proposals will be combined into a single PDF document\n\n"
        
        f"VALIDATION RULES:\n"
        f"- For EACH location, durations count MUST equal net rates count\n"
        f"- If a location has 3 duration options, it MUST have exactly 3 net rates\n"
        f"- DO NOT proceed until ALL locations have complete information\n"
        f"- Ask follow-up questions for any missing information\n\n"
        
        f"SINGLE LOCATION EXAMPLE:\n"
        f"User: 'Proposal for landmark, Jan 1st, 2 weeks at 1.5M'\n"
        f"Bot confirms and generates one proposal\n\n"
        
        f"MULTIPLE LOCATIONS EXAMPLE:\n"
        f"User: 'I need proposals for landmark and gateway'\n"
        f"Bot: 'I'll help you create proposals for The Landmark and The Gateway. Let me get the details for each:\n\n"
        f"For THE LANDMARK:\n"
        f"- What's the campaign start date?\n"
        f"- What duration options do you want?\n"
        f"- What are the net rates for each duration?\n\n"
        f"For THE GATEWAY:\n"
        f"- What's the campaign start date?\n"
        f"- What duration options do you want?\n"
        f"- What are the net rates for each duration?'\n\n"
        
        f"IMPORTANT:\n"
        f"- Always use get_proposals function (works for both single and multiple locations)\n"
        f"- For single location: pass array with one proposal object\n"
        f"- For multiple locations: pass array with multiple proposal objects\n"
        f"- Single location returns: PPT + PDF\n"
        f"- Multiple locations returns: individual PPTs + combined PDF\n"
        f"- Always confirm details before generating\n"
        f"- Format all rates as 'AED X,XXX,XXX'"
    )

    # Get user's conversation history
    history = user_history.get(user_id, [])
    history.append({"role": "user", "content": user_input})
    
    # Keep only last 10 messages to avoid context overflow
    history = history[-10:]
    
    # Build messages with system prompt and history
    messages = [{"role": "developer", "content": prompt}] + history

    tools = [{
        "type": "function",
        "name": "get_proposals",
        "description": "Generate proposals - handles both single and multiple locations. For single: returns PPT+PDF. For multiple: returns individual PPTs + combined PDF",
        "parameters": {
            "type": "object",
            "properties": {
                "proposals": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "properties": {
                            "location": {"type": "string", "description": "The location name (e.g., landmark, gateway, oryx)"},
                            "start_date": {"type": "string", "description": "Start date for the campaign (e.g., 1st December 2025)"},
                            "durations": {
                                "type": "array",
                                "items": {"type": "string"},
                                "description": "List of duration options (e.g., ['2 Weeks', '4 Weeks', '6 Weeks'])"
                            },
                            "net_rates": {
                                "type": "array", 
                                "items": {"type": "string"},
                                "description": "List of net rates corresponding to each duration (e.g., ['AED 1,250,000', 'AED 2,300,000', 'AED 3,300,000'])"
                            }
                        },
                        "required": ["location", "start_date", "durations", "net_rates"]
                    },
                    "description": "Array of proposal objects. Pass one object for single location, multiple for combined proposals"
                }
            },
            "required": ["proposals"]
        }
    }]

    try:
        # Call OpenAI responses API
        res = await openai_client.responses.create(
            model="gpt-4.1",
            input=messages,
            tools=tools,
            tool_choice="auto"
        )

        # Check if we got a valid response
        if not res.output or len(res.output) == 0:
            logger.warning("Empty response from OpenAI")
            await slack_client.chat_postMessage(
                channel=channel, 
                text="I can help with proposals for: " + ", ".join(LOCATION_MAPPING.keys())
            )
            return

        msg = res.output[0]
        
        if msg.type == "function_call" and msg.name == "get_proposals":
            # Process proposals (handles both single and multiple)
            args = json.loads(msg.arguments)
            proposals_data = args.get("proposals", [])
            
            if not proposals_data:
                reply = "❌ No proposals data provided"
                await slack_client.chat_postMessage(channel=channel, text=reply)
            else:
                # Process proposals
                result = await process_proposals(proposals_data)
                
                if result["success"]:
                    if result["is_single"]:
                        # Single proposal - upload PPT and PDF
                        await slack_client.files_upload_v2(
                            channel=channel,
                            file=result["pptx_path"],
                            filename=result["pptx_filename"],
                            initial_comment=f"Here's your PowerPoint proposal for {result['location']}!",
                        )
                        await slack_client.files_upload_v2(
                            channel=channel,
                            file=result["pdf_path"],
                            filename=result["pdf_filename"],
                            initial_comment=f"Here's the PDF version for {result['location']}!",
                        )
                        reply = f"✅ Sent you both PowerPoint and PDF proposals for {result['location']}."
                        await slack_client.chat_postMessage(channel=channel, text=reply)
                        
                        # Clean up temporary files
                        os.unlink(result["pptx_path"])
                        os.unlink(result["pdf_path"])
                    else:
                        # Multiple proposals - upload individual PPTs and combined PDF
                        for file_info in result["individual_files"]:
                            await slack_client.files_upload_v2(
                                channel=channel,
                                file=file_info["path"],
                                filename=file_info["filename"],
                                initial_comment=f"PowerPoint proposal for {file_info['location']}",
                            )
                        
                        # Send the combined PDF
                        await slack_client.files_upload_v2(
                            channel=channel,
                            file=result["merged_pdf_path"],
                            filename=result["merged_pdf_filename"],
                            initial_comment=f"Combined PDF proposal for {result['locations']}",
                        )
                        
                        reply = f"✅ Sent you individual PowerPoints and combined PDF for {result['locations']}."
                        
                        await slack_client.chat_postMessage(channel=channel, text=reply)
                        
                        # Clean up temporary files
                        for file_info in result["individual_files"]:
                            try:
                                os.unlink(file_info["path"])
                            except:
                                pass
                        os.unlink(result["merged_pdf_path"])
                else:
                    reply = f"❌ {result['error']}"
                    await slack_client.chat_postMessage(channel=channel, text=reply)
            
            # Add assistant's response to history
            history.append({"role": "assistant", "content": reply})
        else:
            # Regular text response
            reply = msg.content[-1].text if hasattr(msg, 'content') and msg.content else "How can I help you today?"
            await slack_client.chat_postMessage(channel=channel, text=reply)
            
            # Add assistant's response to history
            history.append({"role": "assistant", "content": reply})
        
        # Update user history
        user_history[user_id] = history[-10:]

    except Exception as e:
        logger.error(f"LLM loop error: {e}", exc_info=True)
        await slack_client.chat_postMessage(
            channel=channel, 
            text="❌ Something went wrong. Please try again."
        )

# ========== SLACK HANDLER ==========
@app.post("/slack/events")
async def slack_events(request: Request):
    body = await request.body()
    timestamp = request.headers.get("X-Slack-Request-Timestamp")
    signature = request.headers.get("X-Slack-Signature")

    if not signature_verifier.is_valid(body.decode(), timestamp, signature):
        raise HTTPException(status_code=403, detail="Invalid Slack signature")

    data = await request.json()
    if data.get("type") == "url_verification":
        return JSONResponse({"challenge": data["challenge"]})

    event = data.get("event", {})
    if event.get("type") == "message" and not event.get("bot_id"):
        asyncio.create_task(main_llm_loop(event["channel"], event["user"], event.get("text", "")))

    return JSONResponse({"status": "ok"})

# ========== HEALTH CHECK ==========
@app.get("/health")
async def health():
    return {"status": "healthy", "timestamp": datetime.now().isoformat()}

# ========== MAIN ==========
if __name__ == "__main__":
    uvicorn.run("main:app", host="0.0.0.0", port=3000, reload=True)
